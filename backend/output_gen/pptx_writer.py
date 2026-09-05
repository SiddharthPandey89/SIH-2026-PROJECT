"""
backend/output_gen/pptx_writer.py

Local PPTX document writer for the Sovereign AI Workbench.

This module provides offline presentation generation for Microsoft PowerPoint .pptx files
from structured content. It is a workspace-bound tool that enforces local path
security and requires no external APIs or network access.

Public API:
-----------
    write_pptx(content, path, overwrite=False, **kwargs)

        Generate a .pptx file from structured content.

        Args:
            content (dict): Structured presentation content with keys:
                - title (str, optional): Presentation title
                - subtitle (str, optional): Presentation subtitle
                - metadata (dict, optional): Presentation properties
                - slides (list, optional): List of slide dicts

            path (str or Path): Output file path. Must resolve within
                                PROJECT_ROOT/data/ or PROJECT_ROOT/sandbox/.

            overwrite (bool): If False (default), reject if path already exists.
                              If True, replace existing file.

            **kwargs: Reserved for future extensions.

        Returns:
            dict: Stable result shape:
                {
                    "success": bool,
                    "path": str | None,
                    "format": "pptx",
                    "message": str | None,
                    "error": str | None,
                }

Security & Constraints:
-----------------------
- Paths are validated and must remain within allowed workspace roots.
- Path traversal (../) and absolute paths outside allowed roots are rejected.
- No files are executed, evaluated, or interpreted as code.
- No environment variables or secrets are accessed.
- The implementation is completely offline and local.

Supported Content Structure:
-----------------------------
    {
        "title": "Presentation Title",
        "subtitle": "Optional subtitle",
        "metadata": {
            "author": "...",
            "subject": "...",
            "company": "..."
        },
        "slides": [
            {
                "title": "Slide Title",
                "content": ["Bullet 1", "Bullet 2"],
            },
            {
                "title": "Paragraph Slide",
                "paragraphs": ["Paragraph 1", "Paragraph 2"],
            },
            {
                "title": "Table Slide",
                "table": {
                    "headers": ["Col 1", "Col 2"],
                    "rows": [["val1", "val2"], ["val3", "val4"]]
                }
            }
        ]
    }

All fields are optional. Missing or None fields are safely ignored.
"""

from __future__ import annotations

import logging
import os
import tempfile
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

__all__ = ["write_pptx"]

logger = logging.getLogger(__name__)

PROJECT_ROOT = Path(__file__).resolve().parents[2]

ALLOWED_ROOTS: List[Path] = [
    (PROJECT_ROOT / "data").resolve(),
    (PROJECT_ROOT / "sandbox").resolve(),
]

# Constraints
_DEFAULT_MAX_FILE_SIZE = 50_000_000  # 50 MB for pptx
_DEFAULT_MAX_SLIDES = 500
_DEFAULT_MAX_ITEMS_PER_SLIDE = 1_000
_DEFAULT_MAX_ROWS_PER_TABLE = 10_000
_DEFAULT_MAX_COLS_PER_TABLE = 100
_DEFAULT_MAX_TEXT_LENGTH = 1_000_000  # characters per field


def _rejected_result(reason: str) -> Dict[str, Any]:
    """Build a stable rejected result dictionary."""
    return {
        "success": False,
        "path": None,
        "format": "pptx",
        "message": None,
        "error": reason,
    }


def _error_result(path: Optional[str], reason: str) -> Dict[str, Any]:
    """Build a stable error result dictionary."""
    return {
        "success": False,
        "path": path,
        "format": "pptx",
        "message": None,
        "error": reason,
    }


def _success_result(path: str) -> Dict[str, Any]:
    """Build a stable success result dictionary."""
    return {
        "success": True,
        "path": path,
        "format": "pptx",
        "message": "PPTX generated successfully.",
        "error": None,
    }


def _resolve_within_allowed(path_value: Any) -> Optional[Path]:
    """
    Resolve a path and ensure it stays within one of the allowed roots.

    Returns the resolved Path if valid, otherwise None.
    """
    if path_value is None:
        return None
    if not isinstance(path_value, (str, Path)):
        return None
    try:
        candidate = Path(path_value).expanduser().resolve()
    except (TypeError, ValueError, OSError):
        return None
    for root in ALLOWED_ROOTS:
        try:
            candidate.relative_to(root)
        except ValueError:
            continue
        return candidate
    return None


def _validate_output_extension(path: Path) -> bool:
    """
    Validate that the output path has a .pptx extension.

    Returns True if valid, False otherwise.
    """
    if not path.suffix:
        return False
    return path.suffix.lower() == ".pptx"


def _coerce_to_string(value: Any) -> str:
    """
    Convert a supported scalar value to text.

    Content validation rejects unsupported types before generation, so this
    helper intentionally avoids silently stringifying arbitrary objects.
    """
    if value is None:
        return ""
    if isinstance(value, str):
        return value
    if isinstance(value, (int, float, bool)):
        return str(value)
    raise TypeError(f"Unsupported value type: {type(value).__name__}")


def _validate_content(content: Any) -> tuple[bool, Optional[str]]:
    """
    Validate the structure of the content dictionary.

    Returns (is_valid, error_message).
    """
    if content is None:
        return False, "Content cannot be None."
    if not isinstance(content, dict):
        return False, "Content must be a dictionary."

    if len(content) == 0:
        return False, "Content dictionary cannot be empty."

    # Validate title if present
    if "title" in content:
        title = content["title"]
        if title is not None and not isinstance(title, str):
            return False, "Title must be a string or None."
        if isinstance(title, str) and len(title) > _DEFAULT_MAX_TEXT_LENGTH:
            return False, f"Title exceeds maximum length ({_DEFAULT_MAX_TEXT_LENGTH} characters)."

    # Validate subtitle if present
    if "subtitle" in content:
        subtitle = content["subtitle"]
        if subtitle is not None and not isinstance(subtitle, str):
            return False, "Subtitle must be a string or None."
        if isinstance(subtitle, str) and len(subtitle) > _DEFAULT_MAX_TEXT_LENGTH:
            return False, f"Subtitle exceeds maximum length ({_DEFAULT_MAX_TEXT_LENGTH} characters)."

    # Validate metadata if present
    if "metadata" in content:
        metadata = content["metadata"]
        if metadata is not None and not isinstance(metadata, dict):
            return False, "Metadata must be a dictionary or None."
        if isinstance(metadata, dict):
            for key, value in metadata.items():
                if not isinstance(key, str):
                    return False, f"Metadata keys must be strings, got {type(key).__name__}."
                if value is not None and not isinstance(value, (str, int, float, bool)):
                    return False, f"Metadata value for '{key}' must be scalar or None."
                if isinstance(value, str) and len(value) > _DEFAULT_MAX_TEXT_LENGTH:
                    return False, f"Metadata value for '{key}' exceeds maximum length."

    # Validate slides if present
    if "slides" in content:
        slides = content["slides"]
        if slides is None:
            return True, None
        if not isinstance(slides, (list, tuple)):
            return False, "Slides must be a list/tuple or None."
        if len(slides) == 0:
            return False, "Slides list cannot be empty."
        if len(slides) > _DEFAULT_MAX_SLIDES:
            return False, f"Too many slides (max {_DEFAULT_MAX_SLIDES})."

        for i, slide in enumerate(slides):
            if not isinstance(slide, dict):
                return False, f"Slide {i} must be a dictionary."

            # Validate slide title
            if "title" in slide:
                slide_title = slide["title"]
                if slide_title is not None and not isinstance(slide_title, str):
                    return False, f"Slide {i}: title must be a string or None."
                if isinstance(slide_title, str) and len(slide_title) > _DEFAULT_MAX_TEXT_LENGTH:
                    return False, f"Slide {i}: title exceeds maximum length."

            # Validate content if present (bullet list)
            if "content" in slide:
                content_items = slide["content"]
                if content_items is not None and not isinstance(content_items, (list, tuple)):
                    return False, f"Slide {i}: 'content' must be a list/tuple or None."
                if content_items is not None:
                    for item_idx, item in enumerate(content_items):
                        if item is None:
                            continue
                        if not isinstance(item, (str, int, float, bool)):
                            return False, (
                                f"Slide {i}: content item {item_idx} must be a "
                                "string, int, float, bool, or None."
                            )
                        if isinstance(item, str) and len(item) > _DEFAULT_MAX_TEXT_LENGTH:
                            return False, (
                                f"Slide {i}: content item {item_idx} exceeds maximum "
                                f"length ({_DEFAULT_MAX_TEXT_LENGTH} characters)."
                            )

            # Validate paragraphs if present
            if "paragraphs" in slide:
                paragraphs = slide["paragraphs"]
                if paragraphs is not None and not isinstance(paragraphs, (list, tuple)):
                    return False, f"Slide {i}: 'paragraphs' must be a list/tuple or None."
                if paragraphs is not None:
                    for item_idx, item in enumerate(paragraphs):
                        if item is None:
                            continue
                        if not isinstance(item, (str, int, float, bool)):
                            return False, (
                                f"Slide {i}: paragraph item {item_idx} must be a "
                                "string, int, float, bool, or None."
                            )
                        if isinstance(item, str) and len(item) > _DEFAULT_MAX_TEXT_LENGTH:
                            return False, (
                                f"Slide {i}: paragraph item {item_idx} exceeds maximum "
                                f"length ({_DEFAULT_MAX_TEXT_LENGTH} characters)."
                            )

            # Validate table if present
            if "table" in slide:
                table = slide["table"]
                if table is not None and not isinstance(table, dict):
                    return False, f"Slide {i}: 'table' must be a dictionary or None."
                if isinstance(table, dict):
                    if "headers" not in table:
                        return False, f"Slide {i}: table must contain 'headers'."
                    if "rows" not in table:
                        return False, f"Slide {i}: table must contain 'rows'."

                    headers = table["headers"]
                    rows = table["rows"]

                    if not isinstance(headers, (list, tuple)):
                        return False, f"Slide {i}: table headers must be a list/tuple."
                    if len(headers) == 0:
                        return False, f"Slide {i}: table headers cannot be empty."
                    if len(headers) > _DEFAULT_MAX_COLS_PER_TABLE:
                        return False, (
                            f"Slide {i}: table has too many columns "
                            f"(max {_DEFAULT_MAX_COLS_PER_TABLE})."
                        )
                    if not isinstance(rows, (list, tuple)):
                        return False, f"Slide {i}: table rows must be a list/tuple."
                    if len(rows) > _DEFAULT_MAX_ROWS_PER_TABLE:
                        return False, (
                            f"Slide {i}: table has too many rows "
                            f"(max {_DEFAULT_MAX_ROWS_PER_TABLE})."
                        )

                    for row_idx, row in enumerate(rows):
                        if not isinstance(row, (list, tuple)):
                            return False, (
                                f"Slide {i}: table row {row_idx} must be a list/tuple."
                            )
                        if len(row) != len(headers):
                            return False, (
                                f"Slide {i}: table row {row_idx} has {len(row)} cells; "
                                f"expected {len(headers)}."
                            )

                    for cell_idx, cell in enumerate(headers):
                        if cell is not None and not isinstance(cell, (str, int, float, bool)):
                            return False, (
                                f"Slide {i}: table header {cell_idx} has unsupported type."
                            )
                        if isinstance(cell, str) and len(cell) > _DEFAULT_MAX_TEXT_LENGTH:
                            return False, (
                                f"Slide {i}: table header {cell_idx} exceeds maximum length."
                            )

                    for row_idx, row in enumerate(rows):
                        for cell_idx, cell in enumerate(row):
                            if cell is not None and not isinstance(cell, (str, int, float, bool)):
                                return False, (
                                    f"Slide {i}: table cell ({row_idx}, {cell_idx}) "
                                    "has unsupported type."
                                )
                            if isinstance(cell, str) and len(cell) > _DEFAULT_MAX_TEXT_LENGTH:
                                return False, (
                                    f"Slide {i}: table cell ({row_idx}, {cell_idx}) "
                                    "exceeds maximum length."
                                )

    return True, None


def _add_table_to_slide(slide, table_spec: Any) -> Optional[str]:
    """
    Add a table to a slide from a structured specification.
    
    Validates that all rows are list/tuple and have exactly num_cols cells.
    
    Returns:
        None if table added successfully
        str (error message) if validation fails and presentation should be rejected
    """
    if table_spec is None:
        return None

    if not isinstance(table_spec, dict):
        return "Table specification must be a dictionary or None."

    if "headers" not in table_spec:
        return "Table must contain 'headers'."

    if "rows" not in table_spec:
        return "Table must contain 'rows'."

    headers = table_spec["headers"]
    rows = table_spec["rows"]

    if not isinstance(headers, (list, tuple)):
        return "Table headers must be a list or tuple."

    if len(headers) == 0:
        return "Table headers cannot be empty."

    if not isinstance(rows, (list, tuple)):
        return "Table rows must be a list or tuple."

    num_cols = len(headers)
    num_rows = len(rows)

    # REJECT if table exceeds row limit (do not truncate)
    if num_rows > _DEFAULT_MAX_ROWS_PER_TABLE:
        return (
            f"Table has {num_rows} rows, which exceeds the maximum "
            f"({_DEFAULT_MAX_ROWS_PER_TABLE}). Table rejected."
        )

    # REJECT if table exceeds column limit (do not truncate)
    if num_cols > _DEFAULT_MAX_COLS_PER_TABLE:
        return (
            f"Table has {num_cols} columns, which exceeds the maximum "
            f"({_DEFAULT_MAX_COLS_PER_TABLE}). Table rejected."
        )

    # Validate all rows: each must be a list or tuple with exactly num_cols cells
    for row_idx, row_data in enumerate(rows):
        if not isinstance(row_data, (list, tuple)):
            return (
                f"Table row {row_idx}: not a list/tuple "
                f"(got {type(row_data).__name__}). "
                f"Each row must be a list or tuple."
            )
        if len(row_data) != num_cols:
            return (
                f"Table row {row_idx}: has {len(row_data)} cells, "
                f"expected {num_cols} (header count). "
                f"All rows must match header cell count."
            )

    # All rows are valid, create table
    if len(rows) == 0:
        # No data rows, don't create an empty table
        return None

    # Add table to slide
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(4.5)

    table_shape = slide.shapes.add_table(1 + num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    # Add header row
    for col_idx, header in enumerate(headers):
        header_text = _coerce_to_string(header)
        cell = table.cell(0, col_idx)
        cell.text = header_text
        # Bold the header
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(12)

    # Add data rows
    for row_idx, row_data in enumerate(rows, start=1):
        for col_idx in range(num_cols):
            cell_value = row_data[col_idx]
            cell_text = _coerce_to_string(cell_value)
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(11)

    return None


def _add_slide(prs: Presentation, slide_spec: Dict[str, Any]) -> Optional[str]:
    """
    Add a slide to the presentation from a structured specification.
    
    Returns:
        None if slide added successfully
        str (error message) if validation fails and presentation should be rejected
    """
    if not isinstance(slide_spec, dict):
        return "Slide specification must be a dictionary."

    # Calculate total content items upfront
    total_items = 0

    if "content" in slide_spec:
        content_items = slide_spec["content"]
        if isinstance(content_items, list):
            for item in content_items:
                item_text = _coerce_to_string(item).strip()
                if item_text:
                    total_items += 1

    if "paragraphs" in slide_spec:
        paragraphs = slide_spec["paragraphs"]
        if isinstance(paragraphs, list):
            for para in paragraphs:
                para_text = _coerce_to_string(para).strip()
                if para_text:
                    total_items += 1

    # Reject if total items exceed limit
    if total_items > _DEFAULT_MAX_ITEMS_PER_SLIDE:
        return (
            f"Slide has {total_items} total content items, "
            f"which exceeds the maximum ({_DEFAULT_MAX_ITEMS_PER_SLIDE}). "
            f"Slide rejected. Do not silently skip items."
        )

    # All content validated, now add slide
    # Use blank slide layout
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # Add slide title if present
    if "title" in slide_spec:
        slide_title = slide_spec["title"]
        title_text = _coerce_to_string(slide_title).strip()
        if title_text:
            # Add a text box for the title
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(9)
            height = Inches(1)
            title_box = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            p = title_frame.paragraphs[0]
            p.text = title_text
            p.font.size = Pt(40)
            p.font.bold = True

    # Add bullet content if present
    if "content" in slide_spec:
        content_items = slide_spec["content"]
        if isinstance(content_items, list) and len(content_items) > 0:
            left = Inches(0.5)
            top = Inches(1.8)
            width = Inches(9)
            height = Inches(5)
            content_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = content_box.text_frame
            text_frame.word_wrap = True

            for item_idx, item in enumerate(content_items):
                item_text = _coerce_to_string(item).strip()
                if item_text:
                    if item_idx == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    p.text = item_text
                    p.font.size = Pt(18)
                    p.level = 0

    # Add paragraphs if present
    if "paragraphs" in slide_spec:
        paragraphs = slide_spec["paragraphs"]
        if isinstance(paragraphs, list) and len(paragraphs) > 0:
            left = Inches(0.5)
            top = Inches(1.8)
            width = Inches(9)
            height = Inches(5)
            para_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = para_box.text_frame
            text_frame.word_wrap = True

            for para_idx, para in enumerate(paragraphs):
                para_text = _coerce_to_string(para).strip()
                if para_text:
                    if para_idx == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    p.text = para_text
                    p.font.size = Pt(14)

    # Add table if present
    if "table" in slide_spec:
        table_spec = slide_spec["table"]
        table_error = _add_table_to_slide(slide, table_spec)
        if table_error:
            return table_error

    return None


def _add_title_slide(prs: Presentation, title: Optional[str], subtitle: Optional[str]) -> None:
    """Add a title slide to the presentation if title is provided."""
    if title is None:
        return

    title_text = _coerce_to_string(title).strip()
    if not title_text:
        return

    # Use title slide layout
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title_text

    # Set subtitle if provided
    if subtitle is not None:
        subtitle_text = _coerce_to_string(subtitle).strip()
        if subtitle_text:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle_text


def _add_metadata(prs: Presentation, metadata: Optional[Dict[str, Any]]) -> None:
    """Add presentation properties (metadata) if provided."""
    if metadata is None or not isinstance(metadata, dict):
        return

    props = prs.core_properties

    if "author" in metadata:
        author = _coerce_to_string(metadata["author"]).strip()
        if author:
            props.author = author

    if "subject" in metadata:
        subject = _coerce_to_string(metadata["subject"]).strip()
        if subject:
            props.subject = subject

    if "company" in metadata:
        company = _coerce_to_string(metadata["company"]).strip()
        if company:
            props.company = company


def write_pptx(
    content: Any,
    path: Any,
    overwrite: bool = False,
    **kwargs: Any,
) -> Dict[str, Any]:
    """
    Generate a Microsoft PowerPoint .pptx file from structured content.

    Args:
        content (dict): Structured presentation content.
                        See module docstring for schema.
        path (str or Path): Output file path. Must resolve within
                            PROJECT_ROOT/data/ or PROJECT_ROOT/sandbox/.
        overwrite (bool): If False (default), reject if file exists.
                          If True, replace existing file. Default: False.
        **kwargs: Reserved for future use.

    Returns:
        dict: Result dictionary with keys:
            - success (bool): True if file was written successfully.
            - path (str or None): Resolved path if successful.
            - format (str): Always "pptx".
            - message (str or None): Success message if successful.
            - error (str or None): Error message if unsuccessful.
    """

    # Validate and resolve path
    resolved_path = _resolve_within_allowed(path)
    if resolved_path is None:
        return _rejected_result(
            "Invalid 'path': must be inside data/ or sandbox/."
        )

    # Validate file extension
    if not _validate_output_extension(resolved_path):
        return _rejected_result(
            "Invalid file extension: must end with .pptx"
        )

    # Check if file exists and overwrite policy
    if resolved_path.exists():
        if not overwrite:
            return _rejected_result(
                f"File already exists: {resolved_path}. "
                "Set overwrite=True to replace."
            )
        if resolved_path.is_dir():
            return _error_result(
                str(resolved_path),
                "Path is a directory, not a file.",
            )

    # Validate content structure
    is_valid, validation_error = _validate_content(content)
    if not is_valid:
        return _rejected_result(validation_error or "Invalid content.")

    try:
        # Create parent directories if needed
        resolved_path.parent.mkdir(parents=True, exist_ok=True)

        # Create presentation
        prs = Presentation()

        # Add title slide if title is present
        title = content.get("title")
        subtitle = content.get("subtitle")
        if title is not None:
            _add_title_slide(prs, title, subtitle)

        # Add metadata if present
        metadata = content.get("metadata")
        if metadata is not None:
            _add_metadata(prs, metadata)

        # Add slides if present
        slides = content.get("slides", [])
        if isinstance(slides, (list, tuple)):
            for slide_idx, slide_spec in enumerate(slides):
                slide_error = _add_slide(prs, slide_spec)
                if slide_error:
                    # Reject presentation before replacing any final file.
                    return _rejected_result(
                        f"Slide {slide_idx} validation failed: {slide_error}"
                    )

        # Use temporary file for safe atomic write
        # Write to temp file first, then rename on success
        temp_fd, temp_path = tempfile.mkstemp(
            suffix=".pptx",
            dir=str(resolved_path.parent),
            prefix=".tmp_"
        )
        try:
            # Close the file descriptor
            os.close(temp_fd)

            # Save presentation to temp path
            prs.save(temp_path)

            # Verify temp file was created and has content
            temp_path_obj = Path(temp_path)
            if not temp_path_obj.exists():
                return _error_result(
                    str(resolved_path),
                    "Temporary file was not created.",
                )

            temp_file_size = temp_path_obj.stat().st_size
            if temp_file_size == 0:
                temp_path_obj.unlink(missing_ok=True)
                return _error_result(
                    str(resolved_path),
                    "Temporary file is empty (zero bytes).",
                )

            # ENFORCE FILE SIZE LIMIT: reject if temp file exceeds limit
            if temp_file_size > _DEFAULT_MAX_FILE_SIZE:
                temp_path_obj.unlink(missing_ok=True)
                return _rejected_result(
                    f"Generated file size ({temp_file_size} bytes) "
                    f"exceeds maximum limit ({_DEFAULT_MAX_FILE_SIZE} bytes)."
                )

            # Move temp file to final destination (atomic on same filesystem)
            temp_path_obj.replace(resolved_path)

        except Exception:
            # Clean up temp file on failure
            Path(temp_path).unlink(missing_ok=True)
            raise

        # Verify final file
        if not resolved_path.exists():
            return _error_result(
                str(resolved_path),
                "File was not created.",
            )

        file_size = resolved_path.stat().st_size
        if file_size == 0:
            return _error_result(
                str(resolved_path),
                "File is empty (zero bytes).",
            )

        return _success_result(str(resolved_path))

    except PermissionError as exc:
        return _error_result(
            str(resolved_path),
            f"Permission denied: {exc}",
        )

    except OSError as exc:
        return _error_result(
            str(resolved_path),
            f"File system error: {exc}",
        )

    except Exception as exc:
        logger.exception(f"Unexpected error generating PPTX: {exc}")
        return _error_result(
            str(resolved_path),
            f"Unexpected error: {type(exc).__name__}",
        )
