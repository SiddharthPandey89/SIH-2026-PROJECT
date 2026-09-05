"""
Output generation test suite for the Sovereign AI Workbench.

Covers:
- DOCX writer
- PPTX writer
- XLSX writer
- happy-path generation
- output existence/non-empty checks
- overwrite protection
- path security
- extension validation
- malformed input rejection
- strict row/column validation where applicable
- generated-file integrity checks
- basic output content verification

Run from the repository root with:
    python -m pytest tests/test_output_gen.py -v
"""

from pathlib import Path
import sys
import zipfile

import pytest


# Make imports work when pytest is launched from the repository root.
PROJECT_ROOT = Path(__file__).resolve().parents[1]
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))


from backend.output_gen.docx_writer import write_docx
from backend.output_gen.pptx_writer import write_pptx
from backend.output_gen.xlsx_writer import write_xlsx


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def assert_success(result):
    assert isinstance(result, dict)
    assert result["success"] is True
    assert result["format"] in {"docx", "pptx", "xlsx"}
    assert result["path"]
    assert result["error"] is None


def assert_rejected(result):
    assert isinstance(result, dict)
    assert result["success"] is False
    assert result["error"]


def assert_non_empty_file(path: Path):
    assert path.exists()
    assert path.is_file()
    assert path.stat().st_size > 0


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------

DOCX_CONTENT = {
    "title": "Output Generator Test",
    "metadata": {
        "author": "Test Suite",
        "subject": "DOCX generation test",
    },
    "sections": [
        {
            "heading": "Introduction",
            "paragraphs": [
                "This is a local/offline DOCX generation test.",
                "The output should be readable by python-docx.",
            ],
            "bullets": [
                "Bullet one",
                "Bullet two",
            ],
            "numbers": [
                "Step one",
                "Step two",
            ],
        },
        {
            "heading": "Table",
            "table": {
                "headers": ["Name", "Value"],
                "rows": [
                    ["Alpha", 10],
                    ["Beta", 20],
                ],
            },
        },
    ],
}


def test_docx_happy_path(tmp_path):
    output = tmp_path / "docx_test.docx"

    result = write_docx(DOCX_CONTENT, output)

    assert_success(result)
    assert_non_empty_file(output)

    from docx import Document

    doc = Document(output)
    assert doc.paragraphs
    assert any("Output Generator Test" in p.text for p in doc.paragraphs)
    assert any("Introduction" in p.text for p in doc.paragraphs)
    assert len(doc.tables) == 1
    assert doc.tables[0].cell(1, 0).text == "Alpha"


def test_docx_overwrite_protection(tmp_path):
    output = tmp_path / "existing.docx"
    output.write_bytes(b"original")

    result = write_docx(DOCX_CONTENT, output, overwrite=False)

    assert_rejected(result)
    assert output.read_bytes() == b"original"


def test_docx_overwrite_true(tmp_path):
    output = tmp_path / "existing.docx"
    output.write_bytes(b"original")

    result = write_docx(DOCX_CONTENT, output, overwrite=True)

    assert_success(result)
    assert_non_empty_file(output)


def test_docx_rejects_outside_workspace():
    # A path clearly outside the project workspace must be rejected.
    result = write_docx(DOCX_CONTENT, PROJECT_ROOT.parent / "outside.docx")
    assert_rejected(result)


def test_docx_rejects_wrong_extension(tmp_path):
    result = write_docx(DOCX_CONTENT, tmp_path / "wrong.txt")
    assert_rejected(result)


def test_docx_rejects_malformed_table(tmp_path):
    content = {
        "title": "Bad Table",
        "sections": [
            {
                "heading": "Bad",
                "table": {
                    "headers": ["A", "B"],
                    "rows": [["only-one-cell"]],
                },
            }
        ],
    }

    result = write_docx(content, tmp_path / "bad.docx")

    assert_rejected(result)


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------

PPTX_CONTENT = {
    "title": "Output Generator Test",
    "subtitle": "PPTX generation test",
    "slides": [
        {
            "title": "Introduction",
            "content": [
                "This is a local/offline PPTX generation test.",
                "The presentation should be readable by python-pptx.",
            ],
        },
        {
            "title": "Table",
            "table": {
                "headers": ["Name", "Value"],
                "rows": [
                    ["Alpha", 10],
                    ["Beta", 20],
                ],
            },
        },
    ],
}


def test_pptx_happy_path(tmp_path):
    output = tmp_path / "pptx_test.pptx"

    result = write_pptx(PPTX_CONTENT, output)

    assert_success(result)
    assert_non_empty_file(output)

    from pptx import Presentation

    prs = Presentation(output)
    assert len(prs.slides) == 2
    assert any("Introduction" in shape.text for shape in prs.slides[0].shapes if hasattr(shape, "text"))


def test_pptx_is_valid_zip_package(tmp_path):
    output = tmp_path / "valid.pptx"

    result = write_pptx(PPTX_CONTENT, output)

    assert_success(result)
    assert zipfile.is_zipfile(output)

    with zipfile.ZipFile(output) as zf:
        names = zf.namelist()
        assert "[Content_Types].xml" in names
        assert "ppt/presentation.xml" in names


def test_pptx_overwrite_protection(tmp_path):
    output = tmp_path / "existing.pptx"
    output.write_bytes(b"original")

    result = write_pptx(PPTX_CONTENT, output, overwrite=False)

    assert_rejected(result)
    assert output.read_bytes() == b"original"


def test_pptx_overwrite_true(tmp_path):
    output = tmp_path / "existing.pptx"
    output.write_bytes(b"original")

    result = write_pptx(PPTX_CONTENT, output, overwrite=True)

    assert_success(result)
    assert_non_empty_file(output)


def test_pptx_rejects_outside_workspace():
    result = write_pptx(PPTX_CONTENT, PROJECT_ROOT.parent / "outside.pptx")
    assert_rejected(result)


def test_pptx_rejects_wrong_extension(tmp_path):
    result = write_pptx(PPTX_CONTENT, tmp_path / "wrong.txt")
    assert_rejected(result)


def test_pptx_rejects_malformed_table(tmp_path):
    content = {
        "title": "Bad Table",
        "slides": [
            {
                "title": "Bad",
                "table": {
                    "headers": ["A", "B"],
                    "rows": [["only-one-cell"]],
                },
            }
        ],
    }

    result = write_pptx(content, tmp_path / "bad.pptx")

    assert_rejected(result)


def test_pptx_rejects_too_many_slides(tmp_path):
    content = {
        "title": "Too Many",
        "slides": [
            {"title": f"Slide {i}", "content": ["test"]}
            for i in range(501)
        ],
    }

    result = write_pptx(content, tmp_path / "too_many.pptx")

    assert_rejected(result)


# ---------------------------------------------------------------------------
# XLSX
# ---------------------------------------------------------------------------

XLSX_CONTENT = {
    "workbook": {
        "title": "Output Generator Test",
        "properties": {
            "creator": "Test Suite",
            "subject": "XLSX generation test",
            "description": "Local/offline XLSX test",
        },
    },
    "sheets": [
        {
            "name": "Data",
            "headers": ["Name", "Value", "Active"],
            "rows": [
                ["Alpha", 10, True],
                ["Beta", 20, False],
                ["Gamma", None, True],
            ],
        },
        {
            "name": "EmptyRows",
            "headers": ["Column"],
            "rows": [],
        },
    ],
}


def test_xlsx_happy_path(tmp_path):
    output = tmp_path / "xlsx_test.xlsx"

    result = write_xlsx(XLSX_CONTENT, output)

    assert_success(result)
    assert_non_empty_file(output)

    from openpyxl import load_workbook

    wb = load_workbook(output, read_only=True, data_only=False)

    assert wb.sheetnames == ["Data", "EmptyRows"]
    assert wb["Data"]["A1"].value == "Name"
    assert wb["Data"]["A2"].value == "Alpha"
    assert wb["Data"]["B2"].value == 10
    assert wb["Data"]["C2"].value is True

    wb.close()


def test_xlsx_overwrite_protection(tmp_path):
    output = tmp_path / "existing.xlsx"
    output.write_bytes(b"original")

    result = write_xlsx(XLSX_CONTENT, output, overwrite=False)

    assert_rejected(result)
    assert output.read_bytes() == b"original"


def test_xlsx_overwrite_true(tmp_path):
    output = tmp_path / "existing.xlsx"
    output.write_bytes(b"original")

    result = write_xlsx(XLSX_CONTENT, output, overwrite=True)

    assert_success(result)
    assert_non_empty_file(output)


def test_xlsx_rejects_outside_workspace():
    result = write_xlsx(XLSX_CONTENT, PROJECT_ROOT.parent / "outside.xlsx")
    assert_rejected(result)


def test_xlsx_rejects_wrong_extension(tmp_path):
    result = write_xlsx(XLSX_CONTENT, tmp_path / "wrong.txt")
    assert_rejected(result)


def test_xlsx_rejects_duplicate_sheet_names(tmp_path):
    content = {
        "sheets": [
            {"name": "Data", "headers": ["A"], "rows": [[1]]},
            {"name": "Data", "headers": ["A"], "rows": [[2]]},
        ]
    }

    result = write_xlsx(content, tmp_path / "duplicate.xlsx")

    assert_rejected(result)


def test_xlsx_rejects_invalid_sheet_name(tmp_path):
    content = {
        "sheets": [
            {"name": "Invalid/Name", "headers": ["A"], "rows": [[1]]},
        ]
    }

    result = write_xlsx(content, tmp_path / "invalid_name.xlsx")

    assert_rejected(result)


def test_xlsx_rejects_sheet_name_over_31_chars(tmp_path):
    content = {
        "sheets": [
            {"name": "A" * 32, "headers": ["A"], "rows": [[1]]},
        ]
    }

    result = write_xlsx(content, tmp_path / "long_name.xlsx")

    assert_rejected(result)


def test_xlsx_rejects_duplicate_data_shape_mismatch(tmp_path):
    content = {
        "sheets": [
            {
                "name": "Data",
                "headers": ["A", "B"],
                "rows": [
                    [1, 2],
                    [3],  # mismatch
                ],
            }
        ]
    }

    result = write_xlsx(content, tmp_path / "bad_rows.xlsx")

    assert_rejected(result)


def test_xlsx_rejects_unsupported_cell_type(tmp_path):
    class Unsupported:
        pass

    content = {
        "sheets": [
            {
                "name": "Data",
                "headers": ["A"],
                "rows": [[Unsupported()]],
            }
        ]
    }

    result = write_xlsx(content, tmp_path / "unsupported.xlsx")

    assert_rejected(result)


def test_xlsx_rejects_unsupported_workbook_property(tmp_path):
    content = {
        "workbook": {
            "properties": {
                "creator": "Test",
                "unsupported_property": "must be rejected",
            }
        },
        "sheets": [
            {"name": "Data", "headers": ["A"], "rows": [[1]]},
        ],
    }

    result = write_xlsx(content, tmp_path / "bad_property.xlsx")

    assert_rejected(result)


def test_xlsx_rejects_too_many_columns(tmp_path):
    content = {
        "sheets": [
            {
                "name": "Data",
                "headers": [f"C{i}" for i in range(101)],
                "rows": [],
            }
        ]
    }

    result = write_xlsx(content, tmp_path / "too_many_columns.xlsx")

    assert_rejected(result)


def test_xlsx_rejects_too_many_rows(tmp_path):
    content = {
        "sheets": [
            {
                "name": "Data",
                "headers": ["A"],
                "rows": [[i] for i in range(100_001)],
            }
        ]
    }

    result = write_xlsx(content, tmp_path / "too_many_rows.xlsx")

    assert_rejected(result)


# ---------------------------------------------------------------------------
# Cross-writer API contract checks
# ---------------------------------------------------------------------------

@pytest.mark.parametrize(
    ("writer", "fmt", "suffix"),
    [
        (write_docx, "docx", ".docx"),
        (write_pptx, "pptx", ".pptx"),
        (write_xlsx, "xlsx", ".xlsx"),
    ],
)
def test_result_format_contract(writer, fmt, suffix, tmp_path):
    """Every writer must return the same stable result contract."""
    if fmt == "docx":
        content = DOCX_CONTENT
    elif fmt == "pptx":
        content = PPTX_CONTENT
    else:
        content = XLSX_CONTENT

    output = tmp_path / f"contract{suffix}"
    result = writer(content, output)

    assert isinstance(result, dict)
    assert set(result.keys()) == {"success", "path", "format", "message", "error"}
    assert result["format"] == fmt
    assert result["success"] is True
    assert result["path"] == str(output)
    assert result["message"]
    assert result["error"] is None
